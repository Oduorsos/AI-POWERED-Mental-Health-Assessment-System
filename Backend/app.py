import os
import json
import pickle
import logging
from typing import Optional, List, Dict
from datetime import datetime, timedelta

import requests
import numpy as np
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException, Depends, Body, Header
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from sqlalchemy import create_engine, Column, Integer, String, Text, DateTime, ForeignKey, Float
from sqlalchemy.orm import sessionmaker, declarative_base, relationship
from passlib.context import CryptContext
from jose import JWTError, jwt

try:
    import faiss
except Exception:
    faiss = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None

load_dotenv()

# CONFIG
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
OPENROUTER_BASE_URL = os.getenv("OPENROUTER_BASE_URL", "https://openrouter.ai/api/v1")
DATABASE_URL = os.getenv("DATABASE_URL", "sqlite:///./medisos.db")
FAISS_INDEX_PATH = os.getenv("FAISS_INDEX_PATH", "./faiss_ppt.index")
PPTX_PATH = os.getenv("PPTX_PATH", "/mnt/data/AI-powered Mental health Asssessment System.pptx")
JWT_SECRET = os.getenv("JWT_SECRET", "change_this_super_secret")
JWT_ALGORITHM = "HS256"
ACCESS_TOKEN_EXPIRE_MINUTES = 60*24*7

CHAT_ENDPOINT = f"{OPENROUTER_BASE_URL}/chat/completions"
EMBEDDINGS_ENDPOINT = f"{OPENROUTER_BASE_URL}/embeddings"

# Email env
SMTP_HOST = os.getenv("SMTP_HOST")
SMTP_PORT = int(os.getenv("SMTP_PORT", 587))
SMTP_USER = os.getenv("SMTP_USER")
SMTP_PASS = os.getenv("SMTP_PASS")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("medisos")

# DB setup
Base = declarative_base()
engine = create_engine(DATABASE_URL, connect_args={"check_same_thread": False} if "sqlite" in DATABASE_URL else {})
SessionLocal = sessionmaker(bind=engine, autoflush=False, autocommit=False)

# MODELS
class User(Base):
    __tablename__ = "users"
    id = Column(Integer, primary_key=True, index=True)
    first_name = Column(String(100), nullable=False)
    last_name = Column(String(100), nullable=False)
    email = Column(String(150), unique=True, nullable=False)
    age_group = Column(String(50), nullable=True)
    password_hash = Column(String(255), nullable=False)
    created_at = Column(DateTime, default=datetime.utcnow)

class Psychologist(Base):
    __tablename__ = "psychologists"
    id = Column(Integer, primary_key=True, index=True)
    name = Column(String(200))
    email = Column(String(255))
    phone = Column(String(50), nullable=True)
    notes = Column(Text, nullable=True)

class Session(Base):
    __tablename__ = "sessions"
    id = Column(Integer, primary_key=True, autoincrement=True)  # auto-increment int
    user_id = Column(Integer, ForeignKey("users.id"), nullable=True)
    started_at = Column(DateTime, default=datetime.utcnow)
    ended_at = Column(DateTime, nullable=True)
    status = Column(String(20), default="active")
    created_at = Column(DateTime, default=datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

    user = relationship("User")

class Message(Base):
    __tablename__ = "messages"
    id = Column(Integer, primary_key=True, autoincrement=True)  # auto-increment int
    session_id = Column(Integer, ForeignKey("sessions.id"))
    sender = Column(String(20))  # 'user' or 'assistant'
    text = Column(Text)
    sentiment = Column(Float, nullable=True)
    emotion = Column(String(50), nullable=True)
    risk_score = Column(Integer, nullable=True)
    created_at = Column(DateTime, default=datetime.utcnow)

    session = relationship("Session")

class Report(Base):
    __tablename__ = "reports"
    id = Column(Integer, primary_key=True, autoincrement=True)  # auto-increment int
    user_id = Column(Integer, ForeignKey("users.id"), nullable=True)
    session_id = Column(Integer, ForeignKey("sessions.id"), nullable=True)
    psychologist_id = Column(Integer, nullable=True)
    summary = Column(Text)
    risk_score = Column(Integer)
    urgency = Column(String(20), nullable=True)
    created_at = Column(DateTime, default=datetime.utcnow)
    updated_at = Column(DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

    user = relationship("User")
    session = relationship("Session")

# create tables if missing
Base.metadata.create_all(bind=engine)

# Security helpers
pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
def get_password_hash(password: str) -> str:
    return pwd_context.hash(password)
def verify_password(raw_password: str, hashed: str) -> bool:
    return pwd_context.verify(raw_password, hashed)
def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    expire = datetime.utcnow() + (expires_delta if expires_delta else timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES))
    to_encode.update({"exp": expire})
    return jwt.encode(to_encode, JWT_SECRET, algorithm=JWT_ALGORITHM)
def decode_token(token: str):
    try:
        return jwt.decode(token, JWT_SECRET, algorithms=[JWT_ALGORITHM])
    except JWTError as e:
        raise HTTPException(status_code=401, detail="Invalid token") from e
def get_user_from_token(token: str, db):
    payload = decode_token(token)
    sub = payload.get("sub")
    if not sub:
        raise HTTPException(status_code=401, detail="Invalid token payload")
    user = db.query(User).filter(User.id == int(sub)).first()
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    return user

HEADERS = {"Authorization": f"Bearer {OPENROUTER_API_KEY}", "Content-Type": "application/json"}

def call_openrouter_chat(messages: List[Dict], model="openrouter/auto", max_tokens=512, temperature=0.2):
    if OPENROUTER_API_KEY is None:
        raise HTTPException(status_code=500, detail="OpenRouter API key not configured")
    payload = {"model": model, "messages": messages, "max_tokens": max_tokens, "temperature": temperature}
    r = requests.post(CHAT_ENDPOINT, headers=HEADERS, json=payload, timeout=60)
    try:
        r.raise_for_status()
    except Exception:
        logger.error("OpenRouter chat error: %s", r.text)
        raise HTTPException(status_code=502, detail="LLM provider error")
    return r.json()

def get_openrouter_embeddings(texts: List[str], model="text-embedding-3-large"):
    if OPENROUTER_API_KEY is None:
        raise HTTPException(status_code=500, detail="OpenRouter API key not configured")
    payload = {"model": model, "input": texts}
    r = requests.post(EMBEDDINGS_ENDPOINT, headers=HEADERS, json=payload, timeout=60)
    try:
        r.raise_for_status()
    except Exception:
        logger.error("OpenRouter embeddings error: %s", r.text)
        raise HTTPException(status_code=502, detail="Embeddings provider error")
    return r.json().get("data", [])

FAISS_INDEX = None
DOCS = None
EMBED_DIM = None
def load_faiss_index(index_path: str = FAISS_INDEX_PATH):
    global FAISS_INDEX, DOCS, EMBED_DIM
    if faiss is None:
        logger.warning("faiss not installed")
        return False
    if not os.path.exists(index_path) or not os.path.exists("ppt_docs.pkl"):
        logger.info("Index or docs missing. Call /index_ppt.")
        return False
    FAISS_INDEX = faiss.read_index(index_path)
    with open("ppt_docs.pkl","rb") as f:
        DOCS = pickle.load(f)
    EMBED_DIM = FAISS_INDEX.d
    return True

def query_faiss_topk(vector: List[float], top_k: int = 3):
    if FAISS_INDEX is None:
        ok = load_faiss_index()
        if not ok:
            return []
    arr = np.array([vector]).astype("float32")
    D, I = FAISS_INDEX.search(arr, top_k)
    results = []
    for idx, dist in zip(I[0], D[0]):
        if idx < 0 or idx >= len(DOCS): continue
        results.append({"id": int(idx), "text": DOCS[idx], "distance": float(dist)})
    return results

def pptx_to_text(path: str = PPTX_PATH):
    if Presentation is None:
        raise HTTPException(status_code=500, detail="python-pptx not installed")
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail=f"PPTX not found: {path}")
    prs = Presentation(path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    parts.append(t)
        slides_text.append("\n".join(parts) if parts else f"[slide {i+1} - no text]")
    return slides_text

def index_ppt_to_faiss(ppt_path: str = PPTX_PATH, idx_path: str = FAISS_INDEX_PATH, embedding_model="text-embedding-3-large"):
    slides = pptx_to_text(ppt_path)
    if not slides: raise HTTPException(status_code=400, detail="No slide text")
    emb_resp = get_openrouter_embeddings(slides, model=embedding_model)
    vectors = [np.array(item["embedding"], dtype="float32") for item in emb_resp]
    d = vectors[0].shape[0]
    index = faiss.IndexFlatL2(d)
    index.add(np.stack(vectors))
    faiss.write_index(index, idx_path)
    with open("ppt_docs.pkl","wb") as f: pickle.dump(slides, f)
    load_faiss_index(idx_path)
    return {"indexed": len(slides)}

# Safety helpers (unchanged)
URGENT_KEYWORDS = ["kill myself","I want to die","End my life","suicide","hurt myself","It's over for me","I feel like jumping from the roof","Die","hang myself"]
def quick_urgent_check(text: str) -> bool:
    t = text.lower(); return any(k in t for k in URGENT_KEYWORDS)

def llm_classify_risk(text: str) -> dict:
    system = ("You are a safety classifier. Given a user message return ONLY JSON: "
              '{"risk_score":int,"label":"low|medium|high","reason":"short"}')
    messages = [{"role":"system","content":system}, {"role":"user","content":f'Classify:\n"""\n{text}\n"""'}]
    try:
        resp = call_openrouter_chat(messages, max_tokens=80, temperature=0.0)
        assistant_text = resp["choices"][0]["message"]["content"]
        s = assistant_text.find("{"); e = assistant_text.rfind("}")
        if s != -1 and e != -1 and e > s:
            return json.loads(assistant_text[s:e+1])
    except Exception as e:
        logger.warning("LLM risk classify failed: %s", e)
    return {"risk_score":0,"label":"low","reason":"fallback"}

def extract_trailing_json(text: str):
    text = text.strip()
    last_brace = text.rfind("{")
    if last_brace == -1: return text, {}
    candidate = text[last_brace:]
    try:
        data = json.loads(candidate)
        return text[:last_brace].strip(), data
    except Exception:
        return text, {}

# Email
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
def send_email(to_email: str, subject: str, body: str):
    if not SMTP_HOST or not SMTP_USER or not SMTP_PASS:
        logger.warning("SMTP not configured; skipping email")
        return False
    try:
        msg = MIMEMultipart(); msg["From"]=SMTP_USER; msg["To"]=to_email; msg["Subject"]=subject
        msg.attach(MIMEText(body,"plain"))
        s = smtplib.SMTP(SMTP_HOST, SMTP_PORT); s.starttls(); s.login(SMTP_USER, SMTP_PASS)
        s.sendmail(SMTP_USER, to_email, msg.as_string()); s.quit()
        return True
    except Exception as e:
        logger.error("Send email failed: %s", e); return False

# FastAPI app + schemas
app = FastAPI(title="Mental Health AI Assistant Backend (Int IDs)")
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=True, allow_methods=["*"], allow_headers=["*"])

class RegisterIn(BaseModel):
    first_name: str; last_name: str; email: str; password: str; age_group: Optional[str] = None
class LoginIn(BaseModel):
    email: str; password: str
class ChatIn(BaseModel):
    session_id: Optional[int] = None
    message: str

def get_db():
    db = SessionLocal()
    try: yield db
    finally: db.close()

# Auth
@app.post("/register")
def register(payload: RegisterIn, db=Depends(get_db)):
    if db.query(User).filter(User.email==payload.email).first():
        raise HTTPException(status_code=400, detail="Email exists")
    user = User(first_name=payload.first_name.strip(), last_name=payload.last_name.strip(),
                email=payload.email.strip(), age_group=payload.age_group,
                password_hash=get_password_hash(payload.password))
    db.add(user); db.commit(); db.refresh(user)
    token = create_access_token({"sub": str(user.id)})
    return {"access_token": token, "token_type":"bearer", "user_id": user.id}

@app.post("/login")
def login(payload: LoginIn, db=Depends(get_db)):
    user = db.query(User).filter(User.email==payload.email).first()
    if not user or not verify_password(payload.password, user.password_hash):
        raise HTTPException(status_code=401, detail="Invalid credentials")
    token = create_access_token({"sub": str(user.id)})
    return {"access_token": token, "token_type":"bearer", "user_id": user.id}

# Psychologist endpoints
@app.post("/psychologists")
def add_psychologist(data: Dict = Body(...), db=Depends(get_db)):
    p = Psychologist(name=data.get("name"), email=data.get("email"), phone=data.get("phone"), notes=data.get("notes"))
    db.add(p); db.commit(); db.refresh(p)
    return {"ok": True, "psychologist": {"id": p.id, "name": p.name, "email": p.email}}

@app.get("/psychologists")
def list_psychologists(db=Depends(get_db)):
    ps = db.query(Psychologist).all()
    return {"psychologists": [{"id":p.id,"name":p.name,"email":p.email,"phone":p.phone} for p in ps]}

@app.put("/users/{user_id}/psychologist/{psych_id}")
def assign_psychologist(user_id: int, psych_id: int, db=Depends(get_db)):
    user = db.query(User).get(user_id); psych = db.query(Psychologist).get(psych_id)
    if not user or not psych: raise HTTPException(status_code=404, detail="Not found")
    user.psychologist_id = psych.id; db.commit()
    return {"ok": True}

# Index PPT
@app.post("/index_ppt")
def index_ppt(db=Depends(get_db)):
    if faiss is None: raise HTTPException(status_code=500, detail="faiss not installed")
    res = index_ppt_to_faiss()
    return {"ok": True, "indexed": res["indexed"]}

# Chat endpoint
@app.post("/chat")
def chat(payload: ChatIn = Body(...), authorization: Optional[str] = Header(None), db=Depends(get_db)):
    token = None; user = None
    if authorization:
        try:
            token = authorization.split(" ")[1]; user = get_user_from_token(token, db)
        except Exception:
            user = None

    # create or fetch session
    if payload.session_id:
        session = db.query(Session).get(int(payload.session_id))
        if not session:
            # invalid id -> create new
            session = Session(user_id=(user.id if user else None)); db.add(session); db.commit(); db.refresh(session)
    else:
        session = Session(user_id=(user.id if user else None)); db.add(session); db.commit(); db.refresh(session)

    text = payload.message.strip()
    # save user message
    um = Message(session_id=session.id, sender="user", text=text)
    db.add(um); db.commit(); db.refresh(um)

    # quick keyword check
    if quick_urgent_check(text):
        reply = ("I'm very sorry you're feeling this way. If you are in immediate danger, please contact local emergency services now.")
        am = Message(session_id=session.id, sender="assistant", text=reply, risk_score=100); db.add(am); db.commit()
        return {"session_id": session.id, "reply": reply, "emergency": True, "metadata": {"risk_score":100}}

    # LLM risk classification
    risk = llm_classify_risk(text)
    if risk.get("risk_score",0) >= 70:
        reply = ("I am concerned for your safety. Please contact emergency services. Would you like local resources?")
        am = Message(session_id=session.id, sender="assistant", text=reply, risk_score=risk["risk_score"]); db.add(am); db.commit()
        return {"session_id": session.id, "reply": reply, "emergency": True, "metadata": risk}

    # RAG retrieval
    retrieved = []
    try:
        emb = get_openrouter_embeddings([text])[0]["embedding"]
        retrieved = query_faiss_topk(emb, top_k=3)
    except Exception:
        retrieved = []

    # build system prompt
    user_info = ""
    if user:
        user_info = f"User profile: first_name={user.first_name}, last_name={user.last_name}, age_group={user.age_group}."
        if user.psychologist_id:
            psych = db.query(Psychologist).get(user.psychologist_id)
            if psych: user_info += f" Assigned psychologist: {psych.name} ({psych.email})."

    system_prompt = ("You are Mental Health Assessor... (do NOT diagnose). After reply append JSON: {risk_score,int; emotion,str; confidence,float}." + user_info)
    messages = [{"role":"system","content":system_prompt}]
    if retrieved:
        docs_text = "\n\n---\n\n".join([d["text"] for d in retrieved])
        messages.append({"role":"system","content":f"References:\n{docs_text}"})
    messages.append({"role":"user","content":text})

    resp = call_openrouter_chat(messages, max_tokens=512, temperature=0.3)
    try:
        assistant_full = resp["choices"][0]["message"]["content"]
    except Exception:
        assistant_full = resp.get("choices",[{}])[0].get("text","")
    clean, metadata = extract_trailing_json(assistant_full)

    # save assistant message
    am = Message(session_id=session.id, sender="assistant", text=clean, risk_score=int(metadata.get("risk_score", risk.get("risk_score",0))), emotion=metadata.get("emotion"))
    db.add(am); db.commit(); db.refresh(am)

    # create short report
    rep = Report(user_id=(user.id if user else None), session_id=session.id, summary=(clean[:200] + "..."), risk_score=am.risk_score, psychologist_id=(user.psychologist_id if user else None))
    db.add(rep); db.commit()

    return {"session_id": session.id, "reply": clean, "metadata": metadata, "openrouter_raw": resp}

# End session - summarize and email
@app.post("/end_session")
def end_session(body: Dict = Body(...), authorization: Optional[str] = Header(None), db=Depends(get_db)):
    token = None; user = None
    if authorization:
        try:
            token = authorization.split(" ")[1]; user = get_user_from_token(token, db)
        except Exception:
            user = None
    session_id = body.get("session_id")
    if session_id is None:
        raise HTTPException(status_code=400, detail="session_id required")
    msgs = db.query(Message).filter(Message.session_id==int(session_id)).order_by(Message.created_at.asc()).all()
    convo = "\n\n".join([f"{m.sender}: {m.text}" for m in msgs])
    system = ("You are Mental health Assessment summarizer. Return ONLY JSON: {\"summary\":\"...\",\"risk_score\":int,\"urgency\":\"high|moderate|normal\"}")
    messages = [{"role":"system","content":system}, {"role":"user","content":f"Conversation:\n\n{convo}\n\nReturn JSON."}]
    resp = call_openrouter_chat(messages, max_tokens=400, temperature=0.0)
    assistant_text = resp["choices"][0]["message"]["content"]
    s = assistant_text.find("{"); e = assistant_text.rfind("}")
    summary_json = {}
    if s != -1 and e != -1 and e > s:
        try:
            summary_json = json.loads(assistant_text[s:e+1])
        except Exception:
            summary_json = {"summary": assistant_text, "risk_score":0, "urgency":"normal"}
    else:
        summary_json = {"summary": assistant_text, "risk_score":0, "urgency":"normal"}

    # save report
    rep = Report(user_id=(user.id if user else None), session_id=int(session_id), summary=summary_json.get("summary",""), risk_score=int(summary_json.get("risk_score",0)), psychologist_id=(user.psychologist_id if user else None))
    db.add(rep); db.commit()

    # email psychologist if assigned
    sent = False
    if user and user.psychologist_id:
        psych = db.query(Psychologist).get(user.psychologist_id)
        if psych and psych.email:
            subj = f"Mental Assessment Report for {user.first_name} {user.last_name}"
            body_text = (f"Patient: {user.first_name} {user.last_name}\nEmail: {user.email}\n\nSummary:\n{summary_json.get('summary')}\n\nUrgency: {summary_json.get('urgency')}\nRisk Score: {summary_json.get('risk_score')}\n\nFull conversation:\n{convo}")
            sent = send_email(psych.email, subj, body_text)

    return {"ok": True, "report": summary_json, "email_sent": sent}

# Utility endpoints
@app.get("/sessions/{session_id}/messages")
def get_session_messages(session_id: int, db=Depends(get_db)):
    msgs = db.query(Message).filter(Message.session_id==session_id).order_by(Message.created_at.asc()).all()
    return {"session_id": session_id, "messages":[{"id":m.id,"sender":m.sender,"text":m.text,"risk_score":m.risk_score,"emotion":m.emotion,"created_at":m.created_at.isoformat()} for m in msgs]}

@app.get("/me")
def me(authorization: Optional[str] = Header(None), db=Depends(get_db)):
    if not authorization: raise HTTPException(status_code=401, detail="Authorization required")
    token = authorization.split(" ")[1]; user = get_user_from_token(token, db)
    return {"id": user.id, "first_name": user.first_name, "last_name": user.last_name, "email": user.email, "psychologist_id": user.psychologist_id}

@app.get("/")
def root():
    return {"service":"Backend","endpoints":["/register","/login","/index_ppt","/chat","/end_session","/sessions/{id}/messages"]}

# Startup load index
@app.on_event("startup")
def startup_event():
    if faiss is not None:
        load_faiss_index()

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="0.0.0.0", port=int(os.getenv("PORT",8000)), reload=True)
